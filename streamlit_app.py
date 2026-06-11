from __future__ import annotations

from datetime import datetime
from io import BytesIO

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from audit_engine import (
    IssueRecord,
    SCOPE_LABELS,
    analyze_seo_data,
    load_screaming_frog_csv,
)

st.set_page_config(page_title="SEO Audit Automator", page_icon=":mag:", layout="wide")

LEVEL_COLORS = {
    "Faible": "#2E7D32",
    "Moyenne": "#EF6C00",
    "Haute": "#C62828",
}

TIER_COLORS = {
    "P0": "#C62828",
    "P1": "#EF6C00",
    "P2": "#1565C0",
}


def badge_html(label: str, value: str, colors: dict[str, str]) -> str:
    color = colors.get(value, "#6B7280")
    return (
        f"<span style='font-weight:600'>{label}:</span> "
        f"<span style='display:inline-block;padding:0.2rem 0.6rem;border-radius:999px;"
        f"background:{color};color:white;font-weight:700;font-size:0.8rem;'>{value}</span>"
    )


def to_rgb(level: str) -> RGBColor:
    mapping = {
        "Faible": RGBColor(46, 125, 50),
        "Moyenne": RGBColor(239, 108, 0),
        "Haute": RGBColor(198, 40, 40),
    }
    return mapping.get(level, mapping["Moyenne"])


def build_pptx(audit_results: list[IssueRecord], site_name: str, selected_scope: str) -> BytesIO:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "SEO Audit Automator"
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"Audit: {site_name}\n"
        f"Périmètre: {selected_scope}\n"
        f"Généré le {datetime.now().strftime('%d/%m/%Y %H:%M')}"
    )

    for item in audit_results:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.2), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = f"{item.id} - {item.name}"
        title_tf.paragraphs[0].font.size = Pt(26)
        title_tf.paragraphs[0].font.bold = True

        meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.2), Inches(1.0))
        meta_tf = meta_box.text_frame
        meta_tf.text = (
            f"Catégorie: {item.category} | Scope: {item.scope} | Priorité impact: {item.priority} | Ticket: {item.tier}"
        )
        meta_tf.paragraphs[0].font.size = Pt(12)
        meta_tf.paragraphs[0].font.bold = True
        p_meta = meta_tf.add_paragraph()
        p_meta.text = f"Colonnes requises: {', '.join(item.required_columns)}"
        p_meta.font.size = Pt(11)

        constat_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.4))
        constat_tf = constat_box.text_frame
        constat_tf.text = "Constat"
        constat_tf.paragraphs[0].font.size = Pt(20)
        constat_tf.paragraphs[0].font.bold = True
        p1 = constat_tf.add_paragraph()
        p1.text = item.constat
        p1.font.size = Pt(14)
        p1 = constat_tf.add_paragraph()
        p1.text = item.explication_seo
        p1.font.size = Pt(12)

        reco_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.4))
        reco_tf = reco_box.text_frame
        reco_tf.text = "Recommandation"
        reco_tf.paragraphs[0].font.size = Pt(20)
        reco_tf.paragraphs[0].font.bold = True
        p2 = reco_tf.add_paragraph()
        p2.text = item.recommandation
        p2.font.size = Pt(14)

        if item.examples:
            p3 = reco_tf.add_paragraph()
            p3.text = "Exemples:"
            p3.font.bold = True
            p3.font.size = Pt(12)
            for example in item.examples[:3]:
                p_example = reco_tf.add_paragraph()
                p_example.text = f"- {example}"
                p_example.font.size = Pt(11)

        priority_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.4),
            Inches(6.4),
            Inches(3.0),
            Inches(0.6),
        )
        priority_shape.fill.solid()
        priority_shape.fill.fore_color.rgb = to_rgb(item.priority)
        priority_shape.line.color.rgb = RGBColor(255, 255, 255)
        priority_tf = priority_shape.text_frame
        priority_tf.text = f"Priorité: {item.priority}"
        priority_tf.paragraphs[0].font.bold = True
        priority_tf.paragraphs[0].font.size = Pt(13)
        priority_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def render_summary_metrics(issues: list[IssueRecord]) -> None:
    detected = sum(1 for item in issues if item.status == "detected")
    ok = sum(1 for item in issues if item.status == "ok")
    not_evaluable = sum(1 for item in issues if item.status == "not_evaluable")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Scénarios détectés", detected)
    with col2:
        st.metric("Scénarios OK", ok)
    with col3:
        st.metric("Scénarios non évaluables", not_evaluable)


def render_issue_cards(issues: list[IssueRecord]) -> None:
    for item in issues:
        st.markdown("<div class='audit-card'>", unsafe_allow_html=True)
        st.markdown(f"### {item.id} · {item.name}")
        c1, c2, c3, c4, c5 = st.columns([1.0, 1.0, 1.0, 1.0, 1.2])
        with c1:
            st.metric("URLs touchées", item.count)
        with c2:
            st.metric("Impact", f"{item.percentage:.1f}%")
        with c3:
            st.markdown(badge_html("Priorité", item.priority, LEVEL_COLORS), unsafe_allow_html=True)
        with c4:
            st.markdown(badge_html("Complexité", item.complexity, LEVEL_COLORS), unsafe_allow_html=True)
        with c5:
            st.markdown(badge_html("Ticket", item.tier, TIER_COLORS), unsafe_allow_html=True)

        st.caption(f"Catégorie: {item.category} · Scope: {item.scope}")
        st.markdown(f"**Constat**\n\n{item.constat}")
        st.markdown(f"**Recommandation**\n\n{item.recommandation}")
        st.markdown(f"**Explication SEO**\n\n{item.explication_seo}")
        st.markdown(f"**Colonnes requises:** `{', '.join(item.required_columns)}`")
        if item.examples:
            st.markdown("**Exemples d'URLs**")
            for example in item.examples:
                st.code(example)
        st.markdown("</div>", unsafe_allow_html=True)


st.title("SEO Audit Automator")
st.caption(
    "Audit technique déterministe sur exports Screaming Frog FR/EN. "
    "Périmètre par défaut: URLs HTML indexables."
)

st.markdown(
    """
<style>
.audit-card {
    border: 1px solid #e5e7eb;
    border-radius: 12px;
    padding: 1rem;
    margin-bottom: 1rem;
    background: #fcfcfd;
}
</style>
""",
    unsafe_allow_html=True,
)

selected_scope = st.sidebar.selectbox("Périmètre d'audit", options=SCOPE_LABELS, index=0)
enable_advanced = st.sidebar.toggle("Activer les signaux avancés (P2)", value=False)

uploaded_file = st.file_uploader("Dépose ton CSV Screaming Frog", type=["csv"])
if not uploaded_file:
    st.info("Charge un export CSV Screaming Frog pour lancer l'analyse.")
    st.stop()

file_fingerprint = f"{uploaded_file.name}:{uploaded_file.size}:{selected_scope}:{enable_advanced}"
if st.session_state.get("active_file") != file_fingerprint:
    st.session_state["active_file"] = file_fingerprint
    st.session_state.pop("audit_results", None)
    st.session_state.pop("all_issues", None)

try:
    df = load_screaming_frog_csv(uploaded_file)
except Exception as exc:  # noqa: BLE001
    st.error(f"Impossible de lire le CSV: {exc}")
    st.stop()

st.subheader("Aperçu des données")
st.dataframe(df.head(5), use_container_width=True)

try:
    issues, metadata = analyze_seo_data(
        df,
        scope_label=selected_scope,
        include_advanced=enable_advanced,
    )
except KeyError as exc:
    st.error(f"Analyse impossible ({exc})")
    st.stop()

st.session_state["all_issues"] = issues

st.subheader("Vue d'ensemble")
render_summary_metrics(issues)
st.caption(
    f"Scope sélectionné: {metadata['selected_scope']} · "
    f"URLs concernées: {metadata['selected_scope_urls']} · Lignes crawlées: {metadata['total_rows']}"
)

overview_rows = []
for item in issues:
    overview_rows.append(
        {
            "ID": item.id,
            "Catégorie": item.category,
            "Scénario": item.name,
            "Ticket": item.tier,
            "Statut": item.status,
            "URLs": item.count,
            "Impact %": round(item.percentage, 2),
            "Priorité": item.priority,
            "Complexité": item.complexity,
            "Scope": item.scope,
        }
    )
st.dataframe(pd.DataFrame(overview_rows), use_container_width=True, hide_index=True)

not_evaluable = [item for item in issues if item.status == "not_evaluable"]
if not_evaluable:
    st.subheader("Scénarios non évaluables")
    missing_rows = []
    for item in not_evaluable:
        missing_rows.append(
            {
                "ID": item.id,
                "Scénario": item.name,
                "Catégorie": item.category,
                "Colonnes requises": ", ".join(item.required_columns),
                "Raison": item.reason,
            }
        )
    st.dataframe(pd.DataFrame(missing_rows), use_container_width=True, hide_index=True)

if st.button("Générer l'Audit", type="primary"):
    st.session_state["audit_results"] = [item for item in issues if item.status == "detected"]

audit_results: list[IssueRecord] | None = st.session_state.get("audit_results")
if audit_results is not None:
    st.subheader("Résultat de l'audit")
    if not audit_results:
        st.info("Aucun scénario détecté sur la portée sélectionnée.")
    else:
        categories = sorted({item.category for item in audit_results})
        selected_category = st.selectbox("Filtrer par catégorie", ["Toutes"] + categories, index=0)
        if selected_category == "Toutes":
            filtered_results = audit_results
        else:
            filtered_results = [item for item in audit_results if item.category == selected_category]

        render_issue_cards(filtered_results)

        pptx_buffer = build_pptx(
            audit_results=audit_results,
            site_name=uploaded_file.name.replace(".csv", ""),
            selected_scope=selected_scope,
        )
        st.download_button(
            "Télécharger le PPTX",
            data=pptx_buffer.getvalue(),
            file_name=f"audit_seo_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

